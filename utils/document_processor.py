# utils/document_processor.py
import os
import tempfile
import io
from typing import Dict, Any, Optional, List
import logging
import chardet

# PDF processing
import PyPDF2
import pdfplumber

# Office documents
from docx import Document
import openpyxl
import xlrd
from pptx import Presentation

# Data files
import pandas as pd

# Email and web
import email
from email import policy
from bs4 import BeautifulSoup
from config.settings import settings

logger = logging.getLogger(__name__)


class EnhancedDocumentProcessor:
    def __init__(self):
        self.supported_extensions = settings.SUPPORTED_EXTENSIONS

    def is_supported_file(self, filename: str) -> bool:
        """Check if file type is supported."""
        _, ext = os.path.splitext(filename.lower())
        return ext in self.supported_extensions

    def detect_encoding(self, file_content: bytes) -> str:
        """Detect file encoding."""
        try:
            detected = chardet.detect(file_content)
            encoding = detected.get("encoding", "utf-8")
            if encoding is None or detected.get("confidence", 0) < 0.7:
                encoding = "utf-8"
            return encoding
        except Exception:
            return "utf-8"

    def extract_text_from_pdf(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text from PDF files."""
        text_content = ""

        try:
            # Try with pdfplumber first (better for complex layouts)
            with pdfplumber.open(io.BytesIO(file_content)) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text()
                    if page_text:
                        text_content += f"\n\n--- Page {page_num} ---\n\n"
                        text_content += page_text

            if text_content.strip():
                return text_content
        except Exception as e:
            logger.warning(f"pdfplumber failed for {filename}: {e}")

        try:
            # Fallback to PyPDF2
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
            for page_num, page in enumerate(pdf_reader.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    text_content += f"\n\n--- Page {page_num} ---\n\n"
                    text_content += page_text

            return text_content if text_content.strip() else None

        except Exception as e:
            logger.error(f"Failed to extract text from PDF {filename}: {e}")
            return None

    def extract_text_from_docx(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text from Word documents."""
        try:
            doc = Document(io.BytesIO(file_content))
            text_content = ""

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content += paragraph.text + "\n"

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content += " | ".join(row_text) + "\n"

            return text_content if text_content.strip() else None

        except Exception as e:
            logger.error(f"Failed to extract text from Word document {filename}: {e}")
            return None

    def extract_text_from_pptx(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text from PowerPoint presentations."""
        try:
            prs = Presentation(io.BytesIO(file_content))
            text_content = ""

            for slide_num, slide in enumerate(prs.slides, 1):
                text_content += f"\n\n--- Slide {slide_num} ---\n\n"

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += shape.text + "\n"

            return text_content if text_content.strip() else None

        except Exception as e:
            logger.error(f"Failed to extract text from PowerPoint {filename}: {e}")
            return None

    def extract_text_from_excel(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text from Excel files."""
        try:
            # Try with openpyxl first (for .xlsx)
            if filename.lower().endswith(".xlsx"):
                workbook = openpyxl.load_workbook(
                    io.BytesIO(file_content), read_only=True
                )
                text_content = ""

                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text_content += f"\n\n--- Sheet: {sheet_name} ---\n\n"

                    rows_data = []
                    for row in sheet.iter_rows(values_only=True):
                        row_data = [
                            str(cell) if cell is not None else "" for cell in row
                        ]
                        if any(cell.strip() for cell in row_data):  # Skip empty rows
                            rows_data.append(" | ".join(row_data))

                    text_content += "\n".join(rows_data)

                return text_content if text_content.strip() else None

            else:
                # Use xlrd for .xls files
                workbook = xlrd.open_workbook(file_contents=file_content)
                text_content = ""

                for sheet in workbook.sheets():
                    text_content += f"\n\n--- Sheet: {sheet.name} ---\n\n"

                    for row_idx in range(sheet.nrows):
                        row_data = []
                        for col_idx in range(sheet.ncols):
                            cell_value = sheet.cell_value(row_idx, col_idx)
                            row_data.append(str(cell_value) if cell_value else "")

                        if any(cell.strip() for cell in row_data):
                            text_content += " | ".join(row_data) + "\n"

                return text_content if text_content.strip() else None

        except Exception as e:
            logger.error(f"Failed to extract text from Excel file {filename}: {e}")
            return None

    def extract_text_from_csv(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text from CSV files."""
        try:
            # Detect encoding
            encoding = self.detect_encoding(file_content)

            # Try different separators
            separators = [",", ";", "\t", "|"]

            for sep in separators:
                try:
                    df = pd.read_csv(
                        io.BytesIO(file_content), encoding=encoding, sep=sep, nrows=5
                    )
                    if len(df.columns) > 1:  # Successfully parsed
                        # Read the full file
                        df = pd.read_csv(
                            io.BytesIO(file_content), encoding=encoding, sep=sep
                        )

                        # Convert to text format
                        text_content = f"CSV File: {filename}\n\n"
                        text_content += (
                            f"Columns: {' | '.join(df.columns.tolist())}\n\n"
                        )

                        # Add sample of data
                        for idx, row in df.iterrows():
                            row_text = " | ".join(
                                [
                                    str(val) if pd.notna(val) else ""
                                    for val in row.values
                                ]
                            )
                            text_content += row_text + "\n"

                            # Limit to prevent huge files
                            if idx >= 1000:
                                text_content += f"\n... (truncated, showing first 1000 rows of {len(df)} total)\n"
                                break

                        return text_content

                except Exception:
                    continue

            # If all separators fail, treat as plain text
            return file_content.decode(encoding, errors="ignore")

        except Exception as e:
            logger.error(f"Failed to extract text from CSV {filename}: {e}")
            return None

    def extract_text_from_html(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text from HTML files."""
        try:
            encoding = self.detect_encoding(file_content)
            html_content = file_content.decode(encoding, errors="ignore")

            soup = BeautifulSoup(html_content, "html.parser")

            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()

            # Get text and clean it up
            text = soup.get_text()

            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text_content = " ".join(chunk for chunk in chunks if chunk)

            return text_content if text_content.strip() else None

        except Exception as e:
            logger.error(f"Failed to extract text from HTML {filename}: {e}")
            return None

    def extract_text_from_email(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text from email files using built-in email library."""
        try:
            import email
            from email import policy

            # Parse email content
            msg = email.message_from_bytes(file_content, policy=policy.default)

            text_content = f"Email: {filename}\n\n"

            # Extract headers
            if msg["subject"]:
                text_content += f"Subject: {msg['subject']}\n"
            if msg["from"]:
                text_content += f"From: {msg['from']}\n"
            if msg["to"]:
                text_content += f"To: {msg['to']}\n"
            if msg["date"]:
                text_content += f"Date: {msg['date']}\n"

            text_content += "\n--- Email Body ---\n\n"

            # Extract body content
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        body = part.get_content()
                        if body:
                            text_content += body + "\n"
                    elif part.get_content_type() == "text/html":
                        # Fallback to HTML if no plain text
                        html_body = part.get_content()
                        if html_body and "text/plain" not in [
                            p.get_content_type() for p in msg.walk()
                        ]:
                            soup = BeautifulSoup(html_body, "html.parser")
                            text_content += soup.get_text() + "\n"
            else:
                # Single part message
                body = msg.get_content()
                if body:
                    text_content += body + "\n"

            return text_content if text_content.strip() else None

        except Exception as e:
            logger.error(f"Failed to extract text from email {filename}: {e}")
            return None

    def extract_text_from_plain_text(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text from plain text files with better encoding handling."""
        try:
            # Try to detect encoding first
            encoding = self.detect_encoding(file_content)

            # Decode with detected encoding
            text_content = file_content.decode(encoding, errors="ignore")

            # Clean up the text using ftfy for better unicode handling
            try:
                import ftfy

                text_content = ftfy.fix_text(text_content)
            except ImportError:
                # ftfy not available, continue without it
                pass

            # Optional: transliterate unicode to ASCII if needed
            try:
                import unidecode

                # Only transliterate if text has lots of non-ASCII characters
                non_ascii_ratio = (
                    sum(1 for c in text_content if ord(c) > 127) / len(text_content)
                    if text_content
                    else 0
                )
                if (
                    non_ascii_ratio > 0.1
                ):  # If more than 10% non-ASCII, consider transliteration
                    logger.info(
                        f"High non-ASCII content detected in {filename}, applying transliteration"
                    )
                    text_content = unidecode.unidecode(text_content)
            except ImportError:
                # unidecode not available, continue without it
                pass

            return text_content if text_content.strip() else None

        except Exception as e:
            logger.error(f"Failed to extract text from plain text file {filename}: {e}")
            return None

    def extract_text_from_file(
        self, file_content: bytes, filename: str
    ) -> Optional[str]:
        """Extract text content from uploaded file based on file type."""
        _, ext = os.path.splitext(filename.lower())

        try:
            # PDF files
            if ext == ".pdf":
                return self.extract_text_from_pdf(file_content, filename)

            # Word documents
            elif ext in [".docx", ".doc"]:
                return self.extract_text_from_docx(file_content, filename)

            # PowerPoint presentations
            elif ext in [".pptx", ".ppt"]:
                return self.extract_text_from_pptx(file_content, filename)

            # Excel files
            elif ext in [".xlsx", ".xls"]:
                return self.extract_text_from_excel(file_content, filename)

            # CSV files
            elif ext in [".csv", ".tsv"]:
                return self.extract_text_from_csv(file_content, filename)

            # HTML files
            elif ext in [".html", ".htm", ".xhtml"]:
                return self.extract_text_from_html(file_content, filename)

            # Email files
            elif ext == ".eml":
                return self.extract_text_from_email(file_content, filename)

            # Plain text files (including code files)
            else:
                return self.extract_text_from_plain_text(file_content, filename)

        except Exception as e:
            logger.error(f"Unexpected error extracting text from {filename}: {e}")
            return None

    def validate_file_size(self, file_size: int, max_size_mb: int = 10) -> bool:
        """Validate file size."""
        max_size_bytes = max_size_mb * 1024 * 1024
        return file_size <= max_size_bytes

    def get_file_type_description(self, filename: str) -> str:
        """Get human-readable file type description."""
        _, ext = os.path.splitext(filename.lower())

        type_descriptions = {
            ".pdf": "PDF Document",
            ".docx": "Word Document",
            ".doc": "Word Document (Legacy)",
            ".pptx": "PowerPoint Presentation",
            ".ppt": "PowerPoint Presentation (Legacy)",
            ".xlsx": "Excel Spreadsheet",
            ".xls": "Excel Spreadsheet (Legacy)",
            ".csv": "CSV Data File",
            ".tsv": "Tab-Separated Values",
            ".html": "HTML Document",
            ".htm": "HTML Document",
            ".eml": "Email Message",
            ".txt": "Text File",
            ".md": "Markdown Document",
            ".py": "Python Script",
            ".js": "JavaScript File",
            ".json": "JSON Data File",
            ".xml": "XML Document",
            ".yaml": "YAML Configuration",
            ".yml": "YAML Configuration",
        }

        return type_descriptions.get(
            ext, f"{ext.upper()} File" if ext else "Unknown File Type"
        )

    def process_uploaded_file(
        self, uploaded_file, max_size_mb: int = 10
    ) -> Dict[str, Any]:
        """Process an uploaded file and return extracted information."""
        result = {
            "success": False,
            "filename": uploaded_file.name,
            "text_content": None,
            "file_size": 0,
            "file_type": None,
            "file_type_description": None,
            "error_message": None,
        }

        try:
            # Get file info
            file_content = uploaded_file.read()
            result["file_size"] = len(file_content)

            # Get file extension and description
            _, ext = os.path.splitext(uploaded_file.name.lower())
            result["file_type"] = ext
            result["file_type_description"] = self.get_file_type_description(
                uploaded_file.name
            )

            # Validate file size
            if not self.validate_file_size(result["file_size"], max_size_mb):
                result["error_message"] = (
                    f"File size ({result['file_size']/1024/1024:.1f} MB) exceeds maximum allowed size ({max_size_mb} MB)"
                )
                return result

            # Check if file type is supported
            if not self.is_supported_file(uploaded_file.name):
                result["error_message"] = (
                    f"File type '{ext}' is not supported. Supported types: {', '.join(sorted(self.supported_extensions))}"
                )
                return result

            # Extract text content
            text_content = self.extract_text_from_file(file_content, uploaded_file.name)

            if text_content is None:
                result["error_message"] = "Failed to extract text content from file"
                return result

            if len(text_content.strip()) == 0:
                result["error_message"] = (
                    "File appears to be empty or contains no extractable text"
                )
                return result

            result["text_content"] = text_content
            result["success"] = True

            logger.info(
                f"Successfully processed {result['file_type_description']}: {uploaded_file.name} ({result['file_size']} bytes)"
            )

        except Exception as e:
            result["error_message"] = f"Error processing file: {str(e)}"
            logger.error(f"Error processing file {uploaded_file.name}: {e}")

        return result
