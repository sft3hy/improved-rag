# document_processor.py

import io
import logging
from typing import List, Tuple, Optional
import pandas as pd
from langchain_core.documents import Document

# Document processing libraries
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Handles processing of various document formats"""
    
    def __init__(self):
        self.supported_types = {
            'text/plain': self.process_text,
            'text/markdown': self.process_text,
            'application/pdf': self.process_pdf,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self.process_docx,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self.process_pptx,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self.process_excel,
            'application/vnd.ms-excel': self.process_excel,
            'text/csv': self.process_csv,
        }
    
    def get_supported_extensions(self) -> List[str]:
        """Return list of supported file extensions"""
        extensions = ['txt', 'md']
        
        if PDF_AVAILABLE:
            extensions.append('pdf')
        if DOCX_AVAILABLE:
            extensions.append('docx')
        if PPTX_AVAILABLE:
            extensions.append('pptx')
        if EXCEL_AVAILABLE:
            extensions.extend(['xlsx', 'xls'])
        
        extensions.append('csv')
        return extensions
    
    def get_missing_dependencies(self) -> List[str]:
        """Return list of missing dependencies for full functionality"""
        missing = []
        if not PDF_AVAILABLE:
            missing.append("PyPDF2 (for PDF support)")
        if not DOCX_AVAILABLE:
            missing.append("python-docx (for Word document support)")
        if not PPTX_AVAILABLE:
            missing.append("python-pptx (for PowerPoint support)")
        if not EXCEL_AVAILABLE:
            missing.append("openpyxl (for Excel support)")
        return missing
    
    def process_document(self, file_content: bytes, filename: str, file_type: str) -> Tuple[str, bool]:
        """
        Process a document based on its type
        
        Args:
            file_content: Raw file content as bytes
            filename: Name of the file
            file_type: MIME type of the file
            
        Returns:
            Tuple of (processed_text, success_flag)
        """
        try:
            if file_type in self.supported_types:
                processor = self.supported_types[file_type]
                return processor(file_content, filename)
            else:
                # Try to infer from filename extension
                extension = filename.lower().split('.')[-1] if '.' in filename else ''
                
                if extension == 'txt' or extension == 'md':
                    return self.process_text(file_content, filename)
                elif extension == 'pdf' and PDF_AVAILABLE:
                    return self.process_pdf(file_content, filename)
                elif extension == 'docx' and DOCX_AVAILABLE:
                    return self.process_docx(file_content, filename)
                elif extension == 'pptx' and PPTX_AVAILABLE:
                    return self.process_pptx(file_content, filename)
                elif extension in ['xlsx', 'xls'] and EXCEL_AVAILABLE:
                    return self.process_excel(file_content, filename)
                elif extension == 'csv':
                    return self.process_csv(file_content, filename)
                else:
                    logger.warning(f"Unsupported file type: {file_type} for file: {filename}")
                    return f"Unsupported file type: {file_type}", False
                    
        except Exception as e:
            logger.error(f"Error processing {filename}: {str(e)}")
            return f"Error processing file: {str(e)}", False
    
    def process_text(self, file_content: bytes, filename: str) -> Tuple[str, bool]:
        """Process plain text and markdown files"""
        try:
            text = file_content.decode('utf-8')
            return text, True
        except UnicodeDecodeError:
            try:
                text = file_content.decode('latin-1')
                logger.warning(f"Used latin-1 encoding for {filename}")
                return text, True
            except Exception as e:
                logger.error(f"Failed to decode {filename}: {str(e)}")
                return f"Failed to decode file: {str(e)}", False
    
    def process_pdf(self, file_content: bytes, filename: str) -> Tuple[str, bool]:
        """Process PDF files"""
        if not PDF_AVAILABLE:
            return "PDF processing requires PyPDF2. Install with: pip install PyPDF2", False
        
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text_content = []
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_content.append(f"--- Page {page_num + 1} ---\n{page_text}")
                except Exception as e:
                    logger.warning(f"Failed to extract text from page {page_num + 1} of {filename}: {str(e)}")
                    continue
            
            if text_content:
                return "\n\n".join(text_content), True
            else:
                return "No text content could be extracted from the PDF", False
                
        except Exception as e:
            logger.error(f"Failed to process PDF {filename}: {str(e)}")
            return f"Failed to process PDF: {str(e)}", False
    
    def process_docx(self, file_content: bytes, filename: str) -> Tuple[str, bool]:
        """Process Word documents"""
        if not DOCX_AVAILABLE:
            return "Word document processing requires python-docx. Install with: pip install python-docx", False
        
        try:
            docx_file = io.BytesIO(file_content)
            doc = DocxDocument(docx_file)
            
            text_content = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract tables
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_text.append(" | ".join(row_text))
                
                if table_text:
                    text_content.append("\n--- Table ---\n" + "\n".join(table_text))
            
            if text_content:
                return "\n\n".join(text_content), True
            else:
                return "No text content could be extracted from the Word document", False
                
        except Exception as e:
            logger.error(f"Failed to process Word document {filename}: {str(e)}")
            return f"Failed to process Word document: {str(e)}", False
    
    def process_pptx(self, file_content: bytes, filename: str) -> Tuple[str, bool]:
        """Process PowerPoint presentations"""
        if not PPTX_AVAILABLE:
            return "PowerPoint processing requires python-pptx. Install with: pip install python-pptx", False
        
        try:
            pptx_file = io.BytesIO(file_content)
            presentation = Presentation(pptx_file)
            
            text_content = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Extract text from tables in slides
                    if shape.shape_type == 19:  # Table type
                        try:
                            table_text = []
                            for row in shape.table.rows:
                                row_text = []
                                for cell in row.cells:
                                    if cell.text.strip():
                                        row_text.append(cell.text.strip())
                                if row_text:
                                    table_text.append(" | ".join(row_text))
                            if table_text:
                                slide_text.append("Table:\n" + "\n".join(table_text))
                        except Exception as e:
                            logger.warning(f"Failed to extract table from slide {slide_num}: {str(e)}")
                
                if len(slide_text) > 1:  # More than just the slide header
                    text_content.append("\n".join(slide_text))
            
            if text_content:
                return "\n\n".join(text_content), True
            else:
                return "No text content could be extracted from the PowerPoint presentation", False
                
        except Exception as e:
            logger.error(f"Failed to process PowerPoint {filename}: {str(e)}")
            return f"Failed to process PowerPoint: {str(e)}", False
    
    def process_excel(self, file_content: bytes, filename: str) -> Tuple[str, bool]:
        """Process Excel files"""
        if not EXCEL_AVAILABLE:
            return "Excel processing requires openpyxl. Install with: pip install openpyxl", False
        
        try:
            excel_file = io.BytesIO(file_content)
            
            # Try to read with pandas first for better handling
            try:
                # Read all sheets
                excel_data = pd.read_excel(excel_file, sheet_name=None, engine='openpyxl')
                
                text_content = []
                for sheet_name, df in excel_data.items():
                    if not df.empty:
                        sheet_text = [f"--- Sheet: {sheet_name} ---"]
                        
                        # Convert DataFrame to string representation
                        df_string = df.to_string(index=False, na_rep='')
                        sheet_text.append(df_string)
                        
                        # Also add a more readable format
                        sheet_text.append("\n--- Summary ---")
                        sheet_text.append(f"Rows: {len(df)}, Columns: {len(df.columns)}")
                        sheet_text.append(f"Columns: {', '.join(df.columns.astype(str))}")
                        
                        text_content.append("\n".join(sheet_text))
                
                if text_content:
                    return "\n\n".join(text_content), True
                else:
                    return "No data found in Excel file", False
                    
            except Exception as pandas_error:
                logger.warning(f"Pandas failed for {filename}, trying openpyxl directly: {str(pandas_error)}")
                
                # Fallback to openpyxl
                from openpyxl import load_workbook
                workbook = load_workbook(excel_file, data_only=True)
                
                text_content = []
                for sheet_name in workbook.sheetnames:
                    worksheet = workbook[sheet_name]
                    sheet_text = [f"--- Sheet: {sheet_name} ---"]
                    
                    for row in worksheet.iter_rows(values_only=True):
                        row_text = [str(cell) if cell is not None else '' for cell in row]
                        if any(cell.strip() for cell in row_text if cell):
                            sheet_text.append(" | ".join(row_text))
                    
                    if len(sheet_text) > 1:
                        text_content.append("\n".join(sheet_text))
                
                if text_content:
                    return "\n\n".join(text_content), True
                else:
                    return "No data found in Excel file", False
                    
        except Exception as e:
            logger.error(f"Failed to process Excel file {filename}: {str(e)}")
            return f"Failed to process Excel file: {str(e)}", False
    
    def process_csv(self, file_content: bytes, filename: str) -> Tuple[str, bool]:
        """Process CSV files"""
        try:
            # Try different encodings
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    csv_text = file_content.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                return "Failed to decode CSV file with any encoding", False
            
            # Try to parse with pandas for better formatting
            try:
                csv_file = io.StringIO(csv_text)
                df = pd.read_csv(csv_file)
                
                text_content = [f"--- CSV File: {filename} ---"]
                text_content.append(f"Rows: {len(df)}, Columns: {len(df.columns)}")
                text_content.append(f"Columns: {', '.join(df.columns.astype(str))}")
                text_content.append("\n--- Data ---")
                text_content.append(df.to_string(index=False, na_rep=''))
                
                return "\n".join(text_content), True
                
            except Exception as pandas_error:
                logger.warning(f"Pandas failed for CSV {filename}, using raw text: {str(pandas_error)}")
                # Fallback to raw CSV text
                return f"--- CSV File: {filename} ---\n{csv_text}", True
                
        except Exception as e:
            logger.error(f"Failed to process CSV file {filename}: {str(e)}")
            return f"Failed to process CSV file: {str(e)}", False